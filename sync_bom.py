"""
Synchronise les onglets composants de main.xlsx vers les PowerPoint sources.
main.xlsx est la source unique de vérité pour les références, quantités et descriptions.

Utilisation :
    python3 sync_bom.py
"""

import shutil
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pathlib import Path
from datetime import datetime

EXCEL_FILE  = Path('main.xlsx')
SOURCES_DIR = Path('sources')

# Mapping PowerPoint → onglet Excel
PPTX_TO_SHEET = {
    'NUM_PAL_PO_B115_BLADE.pptx':                'Blade',
    'NUM_PAL_PO_B115_BLADE_SERVICE.pptx':         'Blade service',
    'NUM_PAL_PO_B115_PRECASTWEB.pptx':            'PCW',
    'NUM_PAL_PO_B115_CARBON_UPPERBEAM_NV.pptx':   'Upper',
    'NUM_PAL_PO_B115_CARBON_LOWERBEAM_NV.pptx':   'Lower',
    'NUM_PAL_PO_B115_WEB.pptx':                   'WEB',
}


def load_bom(sheet_name):
    """
    Charge le BOM d'un onglet de main.xlsx.
    Retourne un dict { référence: {'qty': ..., 'des': ...} }
    """
    df = pd.read_excel(EXCEL_FILE, sheet_name=sheet_name)
    bom = {}
    for _, row in df.iterrows():
        ref = str(row.iloc[0]).strip()
        if not ref or ref == 'nan':
            continue
        qty = row.iloc[1]
        des = str(row.iloc[2]).strip() if pd.notna(row.iloc[2]) else ''
        bom[ref] = {
            'qty': f"{float(qty):.2f}" if pd.notna(qty) else '',
            'des': des,
        }
    return bom


def _qty_differs(a, b):
    """Compare deux valeurs de quantité : numériquement si possible, sinon en texte."""
    try:
        return abs(float(a) - float(b)) > 1e-9
    except (ValueError, TypeError):
        return str(a).strip() != str(b).strip()


def _update_cell(cell, new_text):
    """Met à jour le texte d'une cellule PPT en préservant la taille de police."""
    tf = cell.text_frame
    # Récupérer la taille existante avant d'écraser
    font_size = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        font_size = tf.paragraphs[0].runs[0].font.size

    cell.text = new_text

    # Réappliquer la taille si elle existait
    if font_size and tf.paragraphs:
        for para in tf.paragraphs:
            para.font.size = font_size


def sync_powerpoint(pptx_path, bom, dry_run=False):
    """
    Met à jour QTY et DES dans chaque tableau du PowerPoint.
    Modifie le fichier source directement (in-place) sauf si dry_run=True.

    Retourne (nb_mis_à_jour, refs_absentes_du_bom, refs_absentes_du_ppt)
    """
    prs = Presentation(pptx_path)
    updated        = 0
    absent_du_bom  = []   # Dans PPT mais pas dans main.xlsx
    refs_ppt_vus   = set()
    changes        = []   # Détail des cellules modifiées

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table   = shape.table
            headers = [c.text.strip().upper() for c in table.rows[0].cells]

            # Trouver les colonnes utiles
            ref_col = next((i for i, h in enumerate(headers)
                            if h in ('RÉFÉRENCE', 'REFERENCE', 'A9B')), None)
            qty_col = next((i for i, h in enumerate(headers) if h == 'QTY'), None)
            des_col = next((i for i, h in enumerate(headers) if 'DES' in h), None)

            if ref_col is None:
                continue

            for row_idx in range(1, len(table.rows)):
                row = table.rows[row_idx]
                ref = row.cells[ref_col].text.strip()

                if not ref:
                    continue

                refs_ppt_vus.add(ref)

                if ref not in bom:
                    absent_du_bom.append(ref)
                    continue

                entry = bom[ref]

                if qty_col is not None and entry['qty']:
                    current = row.cells[qty_col].text.strip()
                    if _qty_differs(current, entry['qty']):
                        changes.append({'ref': ref, 'field': 'QTY',
                                        'old': current, 'new': entry['qty']})
                        _update_cell(row.cells[qty_col], entry['qty'])
                        updated += 1

                if des_col is not None and entry['des']:
                    current = row.cells[des_col].text.strip()
                    if current != entry['des']:
                        changes.append({'ref': ref, 'field': 'DES',
                                        'old': current, 'new': entry['des']})
                        _update_cell(row.cells[des_col], entry['des'])
                        updated += 1

    if not dry_run:
        # Sauvegarder l'original avant écrasement
        backup_dir = pptx_path.parent / 'backups'
        backup_dir.mkdir(exist_ok=True)
        shutil.copy2(pptx_path, backup_dir / pptx_path.name)
        prs.save(pptx_path)

    # Références dans main.xlsx mais absentes du PPT (lignes ajoutées au BOM)
    absent_du_ppt = [ref for ref in bom if ref not in refs_ppt_vus]

    return updated, list(set(absent_du_bom)), absent_du_ppt, changes


def _generate_html_report(results, grand_total, dry_run=False):
    """
    Génère un rapport HTML et le sauvegarde dans le dossier courant.
    results : liste de dicts { pptx, sheet, bom_count, updated, absent_bom, absent_ppt, skipped }
    Retourne le chemin (Path) du fichier créé.
    """
    now = datetime.now()
    date_str = now.strftime('%d/%m/%Y à %H:%M')
    prefix = 'simulation_bom_' if dry_run else 'rapport_bom_'
    filename = f"{prefix}{now.strftime('%d-%m-%y_%H%M')}.html"
    report_path = Path(filename)

    # Nettoyer les anciens rapports du même type
    for old in Path('.').glob(f'{prefix}*.html'):
        if old != report_path:
            old.unlink(missing_ok=True)

    def _ref_rows(refs, color):
        if not refs:
            return ''
        items = ''.join(f'<code>{r}</code>' for r in refs)
        more = ''
        if len(refs) > 10:
            items = ''.join(f'<code>{r}</code>' for r in refs[:10])
            more = f'<code style="color:{color};font-style:italic">… et {len(refs)-10} autre(s)</code>'
        return f'<div class="ref-list" style="border-color:{color}33;background:{color}0D">{items}{more}</div>'

    cards_html = ''
    for r in results:
        if r['skipped']:
            cards_html += f'''
            <div class="file-card skipped-card">
                <div class="file-title">⚠️ {r["pptx"]}</div>
                <div class="file-sub">Non trouvé dans sources/ — ignoré</div>
            </div>'''
            continue

        warn_bom = ''
        if r['absent_bom']:
            warn_bom = f'''
            <div class="warn-block">
                <div class="warn-title">⚠️ {len(r["absent_bom"])} référence(s) dans le PPT mais absente(s) de main.xlsx <span class="hint">(à supprimer manuellement)</span></div>
                {_ref_rows(r["absent_bom"], "#E67E22")}
            </div>'''

        warn_ppt = ''
        if r['absent_ppt']:
            warn_ppt = f'''
            <div class="warn-block">
                <div class="warn-title" style="color:#2980B9">ℹ️ {len(r["absent_ppt"])} référence(s) dans main.xlsx mais absente(s) du PPT <span class="hint">(à ajouter manuellement)</span></div>
                {_ref_rows(r["absent_ppt"], "#2980B9")}
            </div>'''

        changes_html = ''
        if r['changes']:
            MAX_ROWS = 50
            rows_html = ''.join(
                f'<tr>'
                f'<td><code>{c["ref"]}</code></td>'
                f'<td><span class="field-badge">{c["field"]}</span></td>'
                f'<td class="old-val">{c["old"]}</td>'
                f'<td class="new-val">{c["new"]}</td>'
                f'</tr>'
                for c in r['changes'][:MAX_ROWS]
            )
            if len(r['changes']) > MAX_ROWS:
                rows_html += (
                    f'<tr><td colspan="4" style="text-align:center;color:#95A5A6;'
                    f'font-style:italic;padding:6px">… et {len(r["changes"]) - MAX_ROWS} autre(s)</td></tr>'
                )
            changes_html = f'''
            <details class="changes-detail">
                <summary>📝 Voir les {len(r["changes"])} modification(s) détaillées</summary>
                <table class="changes-table">
                    <thead><tr><th>Référence</th><th>Champ</th><th>Avant</th><th>Après</th></tr></thead>
                    <tbody>{rows_html}</tbody>
                </table>
            </details>'''

        ok_color = '#00B4B4' if r['updated'] > 0 else '#95A5A6'
        cards_html += f'''
        <div class="file-card">
            <div class="file-header">
                <div>
                    <div class="file-title">📊 {r["pptx"]}</div>
                    <div class="file-sub">Onglet : <strong>{r["sheet"]}</strong> — {r["bom_count"]} référence(s) dans main.xlsx</div>
                </div>
                <div class="badge" style="background:{ok_color}">✓ {r["updated"]} cellule(s) mise(s) à jour</div>
            </div>
            {warn_bom}
            {warn_ppt}
            {changes_html}
        </div>'''

    total_warns = sum(len(r['absent_bom']) + len(r['absent_ppt']) for r in results if not r['skipped'])
    summary_color = '#00B4B4' if total_warns == 0 else '#E67E22'
    summary_icon  = '✅' if total_warns == 0 else '⚠️'
    sim_banner = (
        '<div class="simulation-banner">🔍 MODE SIMULATION — Aucun fichier n\'a été modifié</div>'
        if dry_run else ''
    )
    updated_label = ('cellule(s) qui seraient mises à jour'
                     if dry_run else 'cellule(s) mise(s) à jour au total')

    html = f'''<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="UTF-8">
<title>Rapport BOM — {now.strftime("%d/%m/%Y")}</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
         background: #EFEFF5; padding: 24px; color: #2C3E50; }}
  .header {{ background: #2B2660; color: white; padding: 22px 28px;
             border-radius: 10px; margin-bottom: 20px; }}
  .header h1 {{ font-size: 20px; font-weight: 700; }}
  .header .sub {{ color: #80DADA; font-size: 13px; margin-top: 6px; }}
  .summary {{ background: white; border-radius: 10px; padding: 16px 22px;
              margin-bottom: 20px; border-left: 5px solid {summary_color};
              display: flex; align-items: center; gap: 14px; }}
  .summary .big {{ font-size: 36px; font-weight: 800; color: {summary_color}; }}
  .summary .txt {{ font-size: 14px; color: #607D8B; }}
  .summary .icon {{ font-size: 28px; }}
  .file-card {{ background: white; border-radius: 10px; padding: 16px 20px;
                margin-bottom: 12px; border: 1px solid #D8D6E8; }}
  .file-card.skipped-card {{ opacity: .6; border-style: dashed; }}
  .file-header {{ display: flex; justify-content: space-between;
                  align-items: flex-start; margin-bottom: 10px; gap: 12px; }}
  .file-title {{ font-size: 14px; font-weight: 600; color: #2B2660; }}
  .file-sub {{ font-size: 12px; color: #607D8B; margin-top: 3px; }}
  .badge {{ background: #00B4B4; color: white; font-size: 12px; font-weight: 600;
            padding: 5px 12px; border-radius: 20px; white-space: nowrap; flex-shrink: 0; }}
  .warn-block {{ margin-top: 10px; }}
  .warn-title {{ font-size: 12px; font-weight: 600; color: #E67E22; margin-bottom: 4px; }}
  .hint {{ font-weight: normal; color: #95A5A6; }}
  .ref-list {{ border-radius: 6px; padding: 8px 12px; border: 1px solid; }}
  .ref-list code {{ display: inline-block; font-size: 11px; color: #444;
                    background: #f7f7f7; border-radius: 3px; padding: 1px 6px;
                    margin: 2px 3px 2px 0; }}
  .skipped-card .file-title {{ color: #7F8C8D; }}
  .skipped-card .file-sub {{ color: #AAB7B8; }}
  .simulation-banner {{ background: #EEEDF5; border: 1px solid #C8C5E0; border-radius: 10px;
                        padding: 12px 22px; margin-bottom: 20px;
                        font-size: 14px; font-weight: 600; color: #2B2660; }}
  .changes-detail {{ margin-top: 10px; }}
  .changes-detail summary {{ font-size: 12px; font-weight: 600; color: #2B2660;
                              cursor: pointer; padding: 4px 0; user-select: none; }}
  .changes-detail summary:hover {{ color: #00B4B4; }}
  .changes-table {{ width: 100%; border-collapse: collapse; margin-top: 8px; font-size: 11px; }}
  .changes-table th {{ background: #EEEDF5; color: #2B2660; font-weight: 600;
                       padding: 5px 10px; text-align: left; border-bottom: 2px solid #D8D6E8; }}
  .changes-table td {{ padding: 4px 10px; border-bottom: 1px solid #EEEDF5; color: #444; }}
  .changes-table tr:last-child td {{ border-bottom: none; }}
  .field-badge {{ font-weight: 700; font-size: 10px; background: #EEEDF5;
                  color: #2B2660; padding: 1px 6px; border-radius: 3px; }}
  .old-val {{ color: #C0392B; }}
  .new-val {{ color: #00B4B4; font-weight: 600; }}
  .footer {{ text-align: center; color: #95A5A6; font-size: 11px; margin-top: 28px; }}
</style>
</head>
<body>

<div class="header">
  <h1>Rapport de synchronisation BOM</h1>
  <div class="sub">Blade B115 — Généré le {date_str}</div>
</div>

{sim_banner}

<div class="summary">
  <div class="icon">{summary_icon}</div>
  <div class="big">{grand_total}</div>
  <div class="txt">{updated_label}<br>
    {"Aucune intervention manuelle requise." if total_warns == 0 else
     f"{total_warns} référence(s) nécessitent une intervention manuelle."}</div>
</div>

{cards_html}

<div class="footer">Rapport généré automatiquement — Blade B115 Logistics · Siemens Gamesa Renewable Energy</div>
</body>
</html>'''

    report_path.write_text(html, encoding='utf-8')
    return report_path


def main(dry_run=False):
    print("=" * 70)
    if dry_run:
        print("  SIMULATION BOM — main.xlsx → PowerPoints (aucun fichier modifié)")
    else:
        print("  SYNCHRONISATION BOM — main.xlsx → PowerPoints")
    print("=" * 70)
    print()
    print("  Source de vérité : main.xlsx (onglets composants)")
    print("  Cible            : fichiers PowerPoint dans sources/")
    if dry_run:
        print("  ⚠️  MODE SIMULATION — les fichiers sources ne seront PAS modifiés")
    print()

    grand_total_updates = 0
    results = []

    for pptx_filename, sheet_name in PPTX_TO_SHEET.items():
        pptx_path = SOURCES_DIR / pptx_filename

        if not pptx_path.exists():
            print(f"⚠️  {pptx_filename} — non trouvé dans sources/ (ignoré)")
            results.append({'pptx': pptx_filename, 'sheet': sheet_name,
                            'bom_count': 0, 'updated': 0,
                            'absent_bom': [], 'absent_ppt': [], 'changes': [], 'skipped': True})
            continue

        print(f"{'─' * 70}")
        print(f"📊  {sheet_name}  →  {pptx_filename}")

        bom = load_bom(sheet_name)
        print(f"    {len(bom)} références dans main.xlsx")

        updated, absent_bom, absent_ppt, changes = sync_powerpoint(pptx_path, bom, dry_run=dry_run)

        print(f"    ✓ {updated} cellule(s) mise(s) à jour")

        if absent_bom:
            print(f"    ⚠️  {len(absent_bom)} référence(s) dans le PPT mais absente(s) de main.xlsx :")
            for ref in absent_bom[:5]:
                print(f"       - {ref}")
            if len(absent_bom) > 5:
                print(f"       ... et {len(absent_bom) - 5} autre(s)")

        if absent_ppt:
            print(f"    ⚠️  {len(absent_ppt)} référence(s) dans main.xlsx mais absente(s) du PPT :")
            for ref in absent_ppt[:5]:
                print(f"       - {ref}")
            if len(absent_ppt) > 5:
                print(f"       ... et {len(absent_ppt) - 5} autre(s)")

        results.append({'pptx': pptx_filename, 'sheet': sheet_name,
                        'bom_count': len(bom), 'updated': updated,
                        'absent_bom': absent_bom, 'absent_ppt': absent_ppt,
                        'changes': changes, 'skipped': False})
        grand_total_updates += updated
        print()

    print("=" * 70)
    if dry_run:
        print(f"🔍  Simulation terminée — {grand_total_updates} cellule(s) seraient mises à jour")
    else:
        print(f"✅  Synchronisation terminée — {grand_total_updates} cellule(s) mise(s) à jour au total")
    print("=" * 70)

    # Générer le rapport HTML
    report_path = _generate_html_report(results, grand_total_updates, dry_run=dry_run)
    print(f"\n📄  Rapport généré : {report_path.resolve()}")
    return report_path


if __name__ == "__main__":
    main()
