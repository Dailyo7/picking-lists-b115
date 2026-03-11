"""
Script universel pour ajouter automatiquement les emplacements dans TOUS les PowerPoint
Compatible avec tous les sous-composants : Blade, Blade service, PCW, Upper, Lower, WEB
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import sys

# Mapping entre les fichiers PowerPoint et les noms de composants dans les picking lists
PPTX_TO_COMPONENT = {
    'NUM_PAL_PO_B115_BLADE.pptx': 'Blade',
    'NUM_PAL_PO_B115_BLADE_SERVICE.pptx': 'Blade_service',
    'NUM_PAL_PO_B115_PRECASTWEB.pptx': 'PCW',
    'NUM_PAL_PO_B115_CARBON_UPPERBEAM_NV.pptx': 'Upper',
    'NUM_PAL_PO_B115_CARBON_LOWERBEAM_NV.pptx': 'Lower',
    'NUM_PAL_PO_B115_WEB.pptx': 'WEB'
}

# Labels recherchés dans les zones de texte PO des PowerPoints (ex: "WEB:    ")
COMPONENT_PO_LABELS = {
    'Blade':         ['BLADE'],
    'Blade_service': ['BLADE SERVICE', 'BLADE_SERVICE'],
    'PCW':           ['PCW', 'PRECASTWEB', 'PRECAST WEB'],
    'Upper':         ['UPPER', 'UPPERBEAM'],
    'Lower':         ['LOWER', 'LOWERBEAM'],
    'WEB':           ['WEB'],
}


def _fill_po_box_in_layout(layout, component_name, po_number):
    """Recherche la zone PO (AUTO_SHAPE dans le layout) et y insère le N° PO.
    Structure réelle : AUTO_SHAPE avec texte "WEB \xa0", "Blade n°", "Upper Beam n°"…
    Le N° PO est ajouté au dernier run du premier paragraphe.
    Retourne True si la zone a été trouvée et mise à jour.
    """
    labels = COMPONENT_PO_LABELS.get(component_name, [component_name.upper()])

    for shape in layout.shapes:
        if shape.shape_type != 1 or not shape.has_text_frame:   # AUTO_SHAPE uniquement
            continue
        full_text_upper = shape.text_frame.text.upper()

        for label in labels:
            if label in full_text_upper:
                # Zone PO trouvée — mettre à jour le dernier run
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        last_run = para.runs[-1]
                        base = last_run.text.rstrip('\xa0').rstrip()
                        last_run.text = (base + ' ' + po_number) if base else po_number
                        return True
    return False


def get_locations_from_picking_list(picking_list_path):
    """
    Extrait les emplacements depuis une picking list
    
    Args:
        picking_list_path: Chemin vers la picking list Excel
    
    Returns:
        Dictionnaire {référence: [liste d'emplacements]}
    """
    # header=1 : les en-têtes sont en ligne 2 depuis l'ajout du titre en ligne 1
    df = pd.read_excel(picking_list_path, header=1)
    
    # Créer un dictionnaire référence -> liste d'emplacements
    locations = {}
    for idx, row in df.iterrows():
        ref = row['Référence']
        location = row['Emplacement']
        
        if ref not in locations:
            locations[ref] = []
        
        # Éviter les doublons
        if location not in locations[ref]:
            locations[ref].append(location)
    
    return locations

def update_powerpoint(pptx_path, locations, output_path, component_name=None, po_number=None):
    """
    Met à jour le PowerPoint en ajoutant une colonne "EMPLACEMENT"
    et en renseignant le N° PO si fourni.

    Args:
        pptx_path: Chemin vers le PowerPoint original
        locations: Dictionnaire {référence: [liste d'emplacements]}
        output_path: Chemin de sortie pour le PowerPoint mis à jour
        component_name: Nom du composant (pour trouver la zone PO)
        po_number: N° PO à inscrire dans la zone dédiée (optionnel)
    """
    print(f"📂 Chargement du PowerPoint: {pptx_path.name}")
    prs = Presentation(pptx_path)

    slides_updated = 0
    references_found = 0
    references_not_found = []
    po_boxes_filled = 0

    # Remplir la zone PO une seule fois par layout unique (la zone est dans le layout)
    if po_number and component_name:
        visited_layouts = set()
        for slide in prs.slides:
            layout = slide.slide_layout
            if id(layout) not in visited_layouts:
                visited_layouts.add(id(layout))
                if _fill_po_box_in_layout(layout, component_name, po_number):
                    po_boxes_filled += 1

    # Parcourir toutes les diapositives
    for slide_idx, slide in enumerate(prs.slides, 1):
        # Chercher les tableaux dans la diapo
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                
                # Vérifier si la colonne "EMPLACEMENT" ou "Emplacement" existe déjà
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                
                # Chercher une colonne contenant "emplacement" (insensible à la casse)
                location_col_idx = None
                for idx, header in enumerate(headers):
                    if "emplacement" in header.lower():
                        location_col_idx = idx
                        break
                
                # Si aucune colonne "emplacement" trouvée, utiliser la dernière colonne
                if location_col_idx is None:
                    if "DETAILS SI ECART" in headers:
                        # Remplacer "DETAILS SI ECART"
                        location_col_idx = headers.index("DETAILS SI ECART")
                        table.rows[0].cells[location_col_idx].text = "EMPLACEMENT"
                    else:
                        # Utiliser la dernière colonne
                        location_col_idx = len(headers) - 1
                        table.rows[0].cells[location_col_idx].text = "EMPLACEMENT"
                
                # Parcourir les lignes du tableau (sauf l'en-tête)
                for row_idx in range(1, len(table.rows)):
                    row = table.rows[row_idx]
                    
                    # Trouver l'index de la colonne A9B (référence)
                    if "A9B" in headers:
                        ref_col_idx = headers.index("A9B")
                    else:
                        ref_col_idx = 0  # Par défaut, première colonne
                    
                    # Extraire la référence
                    reference = row.cells[ref_col_idx].text.strip()
                    
                    if reference and reference != "" and reference != "A9B":
                        # Chercher l'emplacement
                        if reference in locations:
                            # Joindre tous les emplacements avec un retour à la ligne
                            location_text = "\n".join(locations[reference])
                            
                            # Mettre à jour la cellule
                            cell = row.cells[location_col_idx]
                            cell.text = location_text
                            
                            # Formatter le texte
                            if cell.text_frame.paragraphs:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(9)
                                    paragraph.alignment = PP_ALIGN.LEFT
                            
                            references_found += 1
                        else:
                            references_not_found.append(reference)
                            
                            # Mettre un indicateur
                            cell = row.cells[location_col_idx]
                            cell.text = "⚠️ NON TROUVÉ"
                            
                            if cell.text_frame.paragraphs:
                                paragraph = cell.text_frame.paragraphs[0]
                                paragraph.font.size = Pt(9)
                
                slides_updated += 1
                break  # Un seul tableau par diapo normalement
    
    # Sauvegarder le PowerPoint mis à jour
    prs.save(output_path)
    
    return {
        'slides_updated': slides_updated,
        'references_found': references_found,
        'references_not_found': list(set(references_not_found)),
        'po_boxes_filled': po_boxes_filled,
    }

def main(po_numbers=None, shared_dir=None):
    """Fonction principale"""
    print("=" * 80)
    print("  MISE À JOUR AUTOMATIQUE DE TOUS LES POWERPOINT")
    print("=" * 80)
    print()
    
    base = Path(shared_dir) if shared_dir else Path('.')

    # Vérifier que le dossier picking_lists existe
    picking_lists_dir = base / 'picking_lists'

    if not picking_lists_dir.exists():
        print("❌ Erreur: Le dossier 'picking_lists' n'existe pas!")
        print("   Assurez-vous d'avoir généré les picking lists d'abord.")
        return

    # Dossier de sortie
    output_dir = base / 'powerpoints_updated'
    output_dir.mkdir(exist_ok=True)

    # Nettoyer les anciens fichiers PW_*
    old_files = list(output_dir.glob('PW_*.pptx'))
    if old_files:
        for f in old_files:
            f.unlink()
        print(f"🗑️  {len(old_files)} ancien(s) fichier(s) supprimé(s)")

    print(f"📁 Dossier de sortie: {output_dir}/")
    print()
    
    total_stats = {
        'total_files': 0,
        'total_slides': 0,
        'total_references': 0,
        'total_not_found': 0
    }
    failed_files = []
    date_str = pd.Timestamp.now().strftime('%d-%m-%y')

    # Normaliser les clés des N° PO (espaces → underscores pour matcher PPTX_TO_COMPONENT)
    if po_numbers:
        po_numbers = {k.replace(' ', '_'): v for k, v in po_numbers.items()}

    # Traiter chaque PowerPoint
    for pptx_filename, component_name in PPTX_TO_COMPONENT.items():
        pptx_path = Path('sources') / pptx_filename  # sources reste local

        if not pptx_path.exists():
            print(f"⚠️  {pptx_filename} - NON TROUVÉ (ignoré)")
            continue

        print("=" * 80)
        print(f"📊 Traitement: {pptx_filename}")
        print(f"   Composant: {component_name}")
        print("=" * 80)

        try:
            # Trouver la picking list correspondante (la plus récente en date de modification)
            # Format actuel: PL_ComponentName_DD-MM-YY.xlsx
            all_pickings = sorted(picking_lists_dir.glob("PL_*.xlsx"), key=lambda p: p.stat().st_mtime, reverse=True)

            # Filtrer pour trouver la bonne picking list
            # Formats supportés :
            #   ancien : PL_Blade_DD-MM-YY.xlsx
            #   nouveau : PL_#0042_Blade_DD-MM-YY.xlsx
            picking_files = []
            for p in all_pickings:
                parts = p.stem.split('_')
                if len(parts) >= 3:
                    # Ignorer PL (parts[0]), la date (parts[-1]) et les numéros (#XXXX)
                    middle = [s for s in parts[1:-1] if not s.startswith('#')]
                    file_component = '_'.join(middle)
                    if file_component == component_name:
                        picking_files.append(p)
                        break  # Prendre la plus récente

            if not picking_files:
                print(f"❌ Aucune picking list trouvée pour {component_name}")
                print(f"   Recherche dans: {picking_lists_dir}/PL_{component_name}_*.xlsx")
                print()
                continue

            picking_list_path = picking_files[0]
            print(f"📋 Picking list: {picking_list_path.name}")

            # Charger les emplacements
            print(f"📥 Extraction des emplacements...")
            locations = get_locations_from_picking_list(picking_list_path)
            print(f"   ✓ {len(locations)} emplacements chargés")

            # Créer le chemin de sortie
            output_path = output_dir / f"PW_{component_name}_{date_str}.pptx"

            # Mettre à jour le PowerPoint
            po_number = (po_numbers or {}).get(component_name)
            print(f"🔧 Mise à jour en cours...")
            stats = update_powerpoint(pptx_path, locations, output_path,
                                      component_name=component_name, po_number=po_number)

            # Afficher le résumé
            print()
            print(f"✅ Résultats:")
            print(f"   • Diapositives mises à jour: {stats['slides_updated']}")
            print(f"   • Références trouvées: {stats['references_found']}")
            print(f"   • Références non trouvées: {len(stats['references_not_found'])}")
            if po_number:
                if stats.get('po_boxes_filled', 0) > 0:
                    print(f"   • N° PO renseigné : {po_number} ✓")
                else:
                    print(f"   ⚠️  Zone PO non trouvée (N° PO saisi : {po_number})")

            if stats['references_not_found']:
                print(f"\n   ⚠️  Références manquantes:")
                for ref in stats['references_not_found'][:5]:
                    print(f"      - {ref}")
                if len(stats['references_not_found']) > 5:
                    print(f"      ... et {len(stats['references_not_found']) - 5} autres")

            print(f"\n   💾 Sauvegardé: {output_path.name}")
            print()

            # Mettre à jour les statistiques totales
            total_stats['total_files'] += 1
            total_stats['total_slides'] += stats['slides_updated']
            total_stats['total_references'] += stats['references_found']
            total_stats['total_not_found'] += len(stats['references_not_found'])

        except Exception as exc:
            print(f"❌ ERREUR sur {pptx_filename}: {exc}")
            print(f"   → Fichier ignoré, traitement des suivants...\n")
            failed_files.append(pptx_filename)
    
    # Rapport final
    print("=" * 80)
    print("📊 RAPPORT GLOBAL")
    print("=" * 80)
    print(f"Fichiers PowerPoint traités: {total_stats['total_files']}")
    print(f"Diapositives mises à jour: {total_stats['total_slides']}")
    print(f"Références trouvées: {total_stats['total_references']}")
    print(f"Références non trouvées: {total_stats['total_not_found']}")
    if failed_files:
        print(f"\n⚠️  Fichiers en erreur ({len(failed_files)}):")
        for f in failed_files:
            print(f"   • {f}")
    print()
    print(f"✅ Tous les PowerPoint mis à jour sont dans: {output_dir}/")
    print("=" * 80)

if __name__ == "__main__":
    main()
